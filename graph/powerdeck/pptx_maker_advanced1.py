##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
import glob
import scipy.misc
import os

OUTPUT_TAG = "SA market"

# new
prs = pptx.Presentation()
# open
#prs_exists = pptx.Presentation("some_presentation.pptx")

# default slide width
prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
prs.slide_height = 5143500

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# blank slide
#slide = prs.slides.add_slide(prs.slide_layouts[6])

# set title
title = slide.shapes.title
title.text = OUTPUT_TAG


# # TOP LEFT (MST)
# pic_left  = int(prs.slide_width * 0.05)
# pic_top   = int(prs.slide_height * 0.1)
# pic_width = int(prs.slide_width * 0.4)
# pic_height = int(prs.slide_height * 0.4)

# TOP RIGHT (Energy Analytics)
#http://www.wemcouncil.org/wp/wemc-tech-blog-2-plotting-netcdf-w  ith-python/

# Bottom right (Weather analytics)

# BOTTOM left (Weather)


# SA price and production
df_all,price=ancil_load.SA_production_price()
# Extract top x % of times
filtered_times=subset_pricefilter(price)



# Reading in weather solar/wind, 
# use cartopyEarth(date) -> image out soalr and MST


# reading in energy, reading in price all in one dollap not so big since not grid point by grid point
# df_all is all of wind




# These are images to put on the page
imagesTL=glob.glob("/mnt/y/Code/Analysis/graph/Graph_figures/VIC_LGA/present/season/size_prod_color_VF_time*.png")
imagesTR=glob.glob("/home/peter/weatheranalytics/presentation/figures/future/VF/corr_prod_color_VF_20*.png")
# need weather files!
imagesB=glob.glob("/home/peter/weatheranalytics/presentation/figures/future/")
# Other globs

# layout
layout={"TITLE":(.05,0.0,.9,.2),"TL":(.05,.1,.4,.4),"TR":(.45,.1,.4,.4),"B":(.05,.6,.9,.4)}
tiles=["TL","TR","B"]

# change this into loop over timeperiod: with multiple globs that we select from.


for g in imagesTL:
    print(g)
    g_head=os.path.basename(g)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # add title
    tb = slide.shapes.add_textbox(1, 1, 1, 1)
    p = tb.text_frame.add_paragraph()
    p.text = g_head
    #p.font.size = pptx.util.Pt(14)


    # # each tile in the layout
    # for t in tiles:

    #     pic_left  = int(prs.slide_width * layout[t][0])
    #     pic_top   = int(prs.slide_height * layout[t][1])
    #     pic_width = int(prs.slide_width * layout[t][2])
    #     pic_height = int(prs.slide_height * layout[t][3])

    #     img = scipy.misc.imread(g)
    #     # specify default height
    #     #pic_height = int(pic_width * img.shape[0] / img.shape[1])
    #     #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
    #     pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)

prs.save("%s.pptx" % OUTPUT_TAG)