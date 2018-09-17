##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
import glob
import scipy.misc
import os
import fnmatch
import ancil_load
import numpy as np


OUTPUT_TAG = "VIC market 98"

# new
prs = pptx.Presentation()
# open
#prs_exists = pptx.Presentation("some_presentation.pptx")

# default slide width
prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
prs.slide_height = 5143500

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# set title
title = slide.shapes.title
title.text = OUTPUT_TAG

# SA price and production
#df_all,price=ancil_load.SA_production_price()

# For victoria
df_all,price=ancil_load.VIC_lilly_production_price()
price=price.iloc[1:]

# Extract top x % of times
# Needs to be same as that determined from powerdeck/main.ipynb
ftime,fprice=ancil_load.subset_pricefilter(price)
filtered_times = fprice.index # Filtered_times are used iteratively to consturct powerdeck
filtered_times = filtered_times[filtered_times.year==2017]
# Reading in weather solar/wind, 
# use cartopyEarth(date) -> image out soalr and MST


# reading in energy, reading in price all in one dollap not so big since not grid point by grid point
# df_all is all of wind

# Weather:  TL
# Analytics: TR
# MST  : B


# old verison of python
#https://stackoverflow.com/questions/2186525/use-a-glob-to-find-files-recursively-in-python
matches = []
for root, dirnames, filenames in os.walk('/mnt/y/Data/Weather/ECMWF/Saved_plots/wind_temp/Victoria_extreme98/2017'):
    for filename in fnmatch.filter(filenames, '*.png'):
        matches.append(os.path.join(root, filename))

# These are images to put on the page
#imagesTL=glob.glob("/mnt/y/Data/Weather/ECMWF/Saved_plots/wind_temp/Victoria_extreme/*")
imagesTL=matches
imagesTR=glob.glob("/mnt/y/Code/Analysis/graph/Graph_figures/EVENTBASED/Victoria_project98/CBN1/*")
# need weather files!
imagesBL=glob.glob("/mnt/y/Code/Analysis/graph/Graph_figures/EVENTBASED/Victoria_project98/MST/*")
# Other globs
imagesBR=glob.glob("/mnt/y/Code/Analysis/graph/Graph_figures/EVENTBASED/Victoria_project98/PCA1/*")


# layout
# Three Quad
#layout={"TITLE":(.05,0.0,.9,.2),"TL":(.05,.1,.4,.4),"TR":(.45,.1,.4,.4),"B":(.05,.6,.9,.4)}

# Four Quad WORKING
layout = {"TITLE":(.05,0.0,.9,.2),"TL":(.05,.1,.4,.4),"TR":(.45,.1,.4,.4),"BL":(.05,.6,.4,.4),"BR":(.45,.6,.4,.4)}
# Four Quad TEST
#layout = {"TITLE":(0.0,0.0,.7,.2),"TL":(0.0,.1,.5,.45),"TR":(.5,.1,.4,.45),"BL":(0.0,.55,.5,.45),"BR":(.5,.55,.5,.45)}


tiles=["TL","TR","BL","BR"]







# change this into loop over timeperiod: with multiple globs that we select from.
for i,ftimes in enumerate(np.array(filtered_times)):

    # pull images from 3 locations such as imagesTL[i], imagesTR[i], imagesB[i]
    # TL = imagesTL[i]
    # TR = imagesTR[i]
    # B  = imagesB[i]
    Tiles={"TL":imagesTL[i],"TR":imagesTR[i],"BL":imagesBL[i],"BR":imagesBR[i]}

    print(ftimes)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # add title
    tb = slide.shapes.add_textbox(1, 1, 1, 1)
    p = tb.text_frame.add_paragraph()
    p.text = "Time:{0}, Price {1} $/MWh".format(ftimes,fprice[i])
        #p.font.size = pptx.util.Pt(14)


    # each tile in the layout
    for t in tiles:

        pic_left  = int(prs.slide_width * layout[t][0])
        pic_top   = int(prs.slide_height * layout[t][1])
        pic_width = int(prs.slide_width * layout[t][2])
        pic_height = int(prs.slide_height * layout[t][3])

        #img = scipy.misc.imread(g)

        img = scipy.misc.imread(Tiles[t])


        # specify default height
        #pic_height = int(pic_width * img.shape[0] / img.shape[1])
        #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
        pic   = slide.shapes.add_picture(Tiles[t], pic_left, pic_top, pic_width, pic_height)

prs.save("%s.pptx" % OUTPUT_TAG)








# Appendix 

# # TOP LEFT (MST)
# pic_left  = int(prs.slide_width * 0.05)
# pic_top   = int(prs.slide_height * 0.1)
# pic_width = int(prs.slide_width * 0.4)
# pic_height = int(prs.slide_height * 0.4)

# TOP RIGHT (Energy Analytics)
#http://www.wemcouncil.org/wp/wemc-tech-blog-2-plotting-netcdf-w  ith-python/

# Bottom right (Weather analytics)

# BOTTOM left (Weather)